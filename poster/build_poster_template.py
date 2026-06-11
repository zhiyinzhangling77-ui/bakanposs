"""
poster/build_poster_template.py
================================

A0-portrait poster template (PowerPoint .pptx) for the Mediterranean ET
analysis A + B story.

Layout (A0 portrait, 841 × 1189 mm = 33.11" × 46.81"):

  ┌──────────────────────────────────────────────────────┐
  │  TITLE  |  Author / Affiliation                      │  ~3.0"
  ├──────────────────────────┬───────────────────────────┤
  │ 1. INTRODUCTION          │  3. RESULTS — Analysis A  │
  │    + Datasets table      │     Universal τ           │
  │                          │     [Fig 4  v31]          │
  │                          │     Key numbers box       │
  │                          ├───────────────────────────┤
  │ 2. METHODS               │  4. RESULTS — Analysis B  │
  │                          │     Satellite blind spot  │
  │                          │     [Fig 4b v32]          │
  │                          │     [Bias scatter]        │
  │                          │     Key numbers box       │
  ├──────────────────────────┴───────────────────────────┤
  │ 5. TAKE-HOME / IMPLICATIONS                          │
  └──────────────────────────────────────────────────────┘

All body text >= 30 pt. Headings 36-44 pt. Title 80 pt.

Usage:
  pip install python-pptx        # if not yet installed
  python3 poster/build_poster_template.py
  → poster/poster_template_A0.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ----------------------------------------------------------------
# Page size — A0 portrait
# ----------------------------------------------------------------
PAGE_W = Inches(33.11)
PAGE_H = Inches(46.81)

# Colours
NAVY = RGBColor(0x1F, 0x3A, 0x5F)        # section bars
TEAL = RGBColor(0x1D, 0x9E, 0x75)        # Tarazona / Analysis A
ORANGE = RGBColor(0xE8, 0x5D, 0x04)      # Oran / Analysis B accents
PURPLE = RGBColor(0x5A, 0x4F, 0xCF)      # bias / MDE
GREY_BG = RGBColor(0xF3, 0xF4, 0xF6)     # body box background
PLACEHOLDER = RGBColor(0xDD, 0xDD, 0xDD) # figure placeholder fill
TEXT_DARK = RGBColor(0x1A, 0x1A, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Font sizes — all body >= 30 pt
TITLE_PT = 80
AUTHOR_PT = 38
SECTION_PT = 44
SUBHEAD_PT = 36
BODY_PT = 30
CAPTION_PT = 28          # only for figure captions (intentionally below 30)


# ----------------------------------------------------------------
# Helpers
# ----------------------------------------------------------------
def add_textbox(slide, left, top, width, height, text,
                  font_size=BODY_PT, bold=False, color=TEXT_DARK,
                  align=PP_ALIGN.LEFT, fill=None, line=None,
                  anchor=MSO_ANCHOR.TOP):
    """Generic text box. `text` may contain newlines; each line is one
    paragraph with identical styling. Bullet lines starting with '• '
    keep the bullet character verbatim."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    tf.vertical_anchor = anchor

    if fill is not None:
        tb.fill.solid()
        tb.fill.fore_color.rgb = fill
    else:
        tb.fill.background()
    if line is not None:
        tb.line.color.rgb = line
        tb.line.width = Pt(1.0)
    else:
        tb.line.fill.background()

    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = ln
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return tb


def add_section_bar(slide, left, top, width, height, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.25)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.size = Pt(SECTION_PT)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"
    return bar


def add_figure_placeholder(slide, left, top, width, height, label):
    """Light grey rectangle marking where to drop a figure later."""
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = PLACEHOLDER
    box.line.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
    box.line.width = Pt(1.5)
    box.line.dash_style = 7   # dashed
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(SUBHEAD_PT)
    run.font.italic = True
    run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    run.font.name = "Calibri"
    return box


# ----------------------------------------------------------------
# Build the poster
# ----------------------------------------------------------------
def build():
    prs = Presentation()
    prs.slide_width  = PAGE_W
    prs.slide_height = PAGE_H
    blank_layout = prs.slide_layouts[6]    # blank
    slide = prs.slides.add_slide(blank_layout)

    # Margins
    M = Inches(0.6)
    GUTTER = Inches(0.6)

    # =============================================================
    # 0. Title strip (full width, 3.0")
    # =============================================================
    title_h = Inches(3.0)
    title_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                              Emu(0), Emu(0),
                                              PAGE_W, title_h)
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = NAVY
    title_box.line.fill.background()

    # title text
    add_textbox(slide, M, Inches(0.25),
                  PAGE_W - 2 * M, Inches(1.6),
                  "A universal 3-day ET recovery timescale across "
                  "Mediterranean drylands —\n"
                  "and the satellite blind spot it reveals at drip-irrigated orchards",
                  font_size=TITLE_PT, bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER)

    # author / affiliation
    add_textbox(slide, M, Inches(1.95),
                  PAGE_W - 2 * M, Inches(1.0),
                  "Shion Nagamine¹    |    "
                  "¹ Affiliation here (lab, university, city)    |    "
                  "Contact: name@example.org",
                  font_size=AUTHOR_PT, color=WHITE,
                  align=PP_ALIGN.CENTER)

    # =============================================================
    # Layout grid — two columns below the title strip
    # =============================================================
    body_top = title_h + Inches(0.4)
    body_bot = PAGE_H - Inches(0.6)         # bottom margin
    bottom_strip_h = Inches(3.4)            # take-home bar at the bottom
    body_h = body_bot - body_top - bottom_strip_h - Inches(0.4)

    col_w = (PAGE_W - 2 * M - GUTTER) / 2
    left_col_x = M
    right_col_x = M + col_w + GUTTER

    # =============================================================
    # LEFT COLUMN — 1. Introduction, 2. Methods
    # =============================================================
    cur_y = body_top

    # ----- Section bar: 1. Introduction -----
    add_section_bar(slide, left_col_x, cur_y, col_w, Inches(1.0),
                       "1. Introduction")
    cur_y += Inches(1.1)

    # Intro text — narrative
    intro_h = Inches(7.0)
    intro_text = (
        "Question.  After each water input — rain or drip irrigation — "
        "ecosystem evapotranspiration (ET) decays exponentially with a "
        "characteristic timescale τ. Is τ a property of the climate "
        "(universal), or of the management (rainfed vs irrigated)?\n"
        "\n"
        "Why it matters.  If τ is universal, the same recovery model "
        "transfers between rainfed and irrigated systems; only the "
        "amplitude scales with management. This separates the climate "
        "signal from the human one.\n"
        "\n"
        "Why a satellite check.  If τ is universal, a 5-km satellite ET "
        "product should reproduce it. Where it fails reveals what the "
        "satellite cannot see."
    )
    add_textbox(slide, left_col_x, cur_y, col_w, intro_h,
                  intro_text, font_size=BODY_PT,
                  fill=GREY_BG, line=RGBColor(0xCC, 0xCC, 0xCC))
    cur_y += intro_h + Inches(0.3)

    # Datasets sub-section
    add_textbox(slide, left_col_x, cur_y, col_w, Inches(0.9),
                  "Datasets — two contrasting EC sites + a satellite",
                  font_size=SUBHEAD_PT, bold=True, color=NAVY)
    cur_y += Inches(1.0)

    datasets_text = (
        "•  Oran (Spain, 38.82°N, 1.86°W) — rainfed winter cereal, "
        "EC tower 2018–2020 (3 yrs).\n"
        "•  Tarazona (Spain, 39.27°N, 1.94°W) — drip-irrigated almond "
        "orchard, EC tower 2020–2024 (5 yrs).\n"
        "•  Meteosat ETv3 (LSA SAF DMET) — 5-km daily ET from MSG/SEVIRI "
        "30-min, 2018–2024."
    )
    add_textbox(slide, left_col_x, cur_y, col_w, Inches(4.5),
                  datasets_text, font_size=BODY_PT,
                  fill=GREY_BG, line=RGBColor(0xCC, 0xCC, 0xCC))
    cur_y += Inches(4.8)

    # ----- Section bar: 2. Methods -----
    add_section_bar(slide, left_col_x, cur_y, col_w, Inches(1.0),
                       "2. Methods")
    cur_y += Inches(1.1)

    methods_text = (
        "Event detection.  A water-input event = day with rain ≥ 3 mm "
        "(Oran) or irrigation ≥ 0.5 mm (Tarazona). Each event window is "
        "4–14 days, ending at the next event.\n"
        "\n"
        "Exponential fit.\n"
        "    LE(d) = LE_∞ + (LE_0 − LE_∞) · exp(−d/τ)\n"
        "Fitted on the median LE per day-since-event (≥ 3 events / day). "
        "LE_∞ fixed at the observed median for d ≥ 7 d.\n"
        "\n"
        "Uncertainty.  Bootstrap (B = 5000, raw-data resampling) gives "
        "the SE and 95 % CI of τ. Boundary or low-R² fits are excluded.\n"
        "\n"
        "MDE.  Minimum detectable effect for the “no difference” claim:\n"
        "    MDE = 1.96 · √(SE₁² + SE₂²)\n"
        "If |Δτ| < MDE, two strata are statistically indistinguishable.\n"
        "\n"
        "Satellite fair comparison.  Cloudy days dropped via n_obs < 36 "
        "of 48 half-hour steps. Freeze days (Ta_min ≤ 0 °C) excluded. "
        "Paired daily comparison only."
    )
    methods_h = body_top + body_h - cur_y
    add_textbox(slide, left_col_x, cur_y, col_w, methods_h,
                  methods_text, font_size=BODY_PT,
                  fill=GREY_BG, line=RGBColor(0xCC, 0xCC, 0xCC))

    # =============================================================
    # RIGHT COLUMN — 3. Results A, 4. Results B
    # =============================================================
    cur_y = body_top

    # ----- Section bar: 3. Results A -----
    add_section_bar(slide, right_col_x, cur_y, col_w, Inches(1.0),
                       "3. Results — Analysis A:  τ is universal")
    cur_y += Inches(1.1)

    resA_lede = (
        "Both rainfed cereal and drip-irrigated almond recover with the "
        "same τ ≈ 3 d. Only the amplitude scales — by 4.5× — with "
        "irrigation."
    )
    add_textbox(slide, right_col_x, cur_y, col_w, Inches(1.8),
                  resA_lede, font_size=BODY_PT, bold=True, color=TEAL)
    cur_y += Inches(2.0)

    figA_h = Inches(11.5)
    add_figure_placeholder(slide, right_col_x, cur_y, col_w, figA_h,
                              "[ Figure 4  —  output_analysis_A_v31/"
                              "fig04_poster_main_v31.pdf ]\n"
                              "Recovery curves (Tarazona, Oran) + "
                              "τ-comparison bars")
    cur_y += figA_h + Inches(0.2)

    add_textbox(slide, right_col_x, cur_y, col_w, Inches(1.0),
                  "Figure 4. EC LE recovery and τ comparison (see "
                  "caption file v31_poster_caption.txt).",
                  font_size=CAPTION_PT, color=RGBColor(0x55, 0x55, 0x55))
    cur_y += Inches(1.1)

    keyA_text = (
        "Key numbers (EC, active season).\n"
        "•  Tarazona (drip almond, n = 41): τ = 3.36 d  [2.44, 4.90]\n"
        "•  Oran (rainfed cereal, n = 10): τ = 2.82 d  [1.83, 5.48]\n"
        "•  |Δτ| = 0.54 d  <  MDE = 2.15 d  → indistinguishable\n"
        "•  Amplitude: 21 vs 95 W/m²  →  4.5× management scaling"
    )
    add_textbox(slide, right_col_x, cur_y, col_w, Inches(5.5),
                  keyA_text, font_size=BODY_PT, color=TEXT_DARK,
                  fill=RGBColor(0xE6, 0xF4, 0xEE), line=TEAL)
    cur_y += Inches(5.8)

    # ----- Section bar: 4. Results B -----
    add_section_bar(slide, right_col_x, cur_y, col_w, Inches(1.0),
                       "4. Results — Analysis B:  Satellite blind spot")
    cur_y += Inches(1.1)

    resB_lede = (
        "At rainfed Oran, Meteosat ETv3 matches the tower (r = 0.82, "
        "bias +1 %). At drip-irrigated Tarazona, the satellite is flat: "
        "r = 0.07, bias −70 %. The entire management amplitude lives "
        "inside the bias — and the bias itself recovers on the "
        "Analysis-A τ timescale."
    )
    add_textbox(slide, right_col_x, cur_y, col_w, Inches(3.4),
                  resB_lede, font_size=BODY_PT, bold=True, color=ORANGE)
    cur_y += Inches(3.6)

    figB_h = Inches(8.5)
    add_figure_placeholder(slide, right_col_x, cur_y, col_w, figB_h,
                              "[ Figure 4b — output_analysis_A_v32/"
                              "fig04b_tarazona_blindspot_v32.pdf ]\n"
                              "Bias recovery curve  +  τ comparison "
                              "(EC vs Sat vs Bias)")
    cur_y += figB_h + Inches(0.2)

    figS_h = Inches(6.5)
    add_figure_placeholder(slide, right_col_x, cur_y, col_w, figS_h,
                              "[ EC vs Meteosat ETv3 scatter — "
                              "output_bias_stats/fig_scatter_ec_vs_metv3.png ]")
    cur_y += figS_h + Inches(0.3)

    keyB_text = (
        "Key numbers (paired daily, qflag-clean, no freeze days).\n"
        "•  Oran rainfed:    r = 0.82,  bias  +1 %,  KGE = 0.61\n"
        "•  Tarazona drip:   r = 0.07,  bias −70 %,  KGE = −0.21\n"
        "•  Bias-recovery fit:  τ_bias = 4.57 d,  amp_bias = 95 W/m²\n"
        "•  Sat-only fit: τ pegged at floor, R² ≈ 0 (no recovery seen)"
    )
    keyB_remaining = body_top + body_h - cur_y
    add_textbox(slide, right_col_x, cur_y, col_w, keyB_remaining,
                  keyB_text, font_size=BODY_PT, color=TEXT_DARK,
                  fill=RGBColor(0xFD, 0xEC, 0xE3), line=ORANGE)

    # =============================================================
    # 5. Take-home strip (full width, bottom)
    # =============================================================
    th_top = PAGE_H - Inches(0.6) - bottom_strip_h
    th_w = PAGE_W - 2 * M

    add_section_bar(slide, M, th_top, th_w, Inches(1.0),
                       "5. Take-home")
    takehome = (
        "•  τ ≈ 3 days is a climate property of Mediterranean drylands — "
        "rainfed and drip-irrigated share the same recovery clock.\n"
        "•  Amplitude is the management signal: ~95 W/m² under drip "
        "irrigation vs ~21 W/m² rainfed (4.5×).\n"
        "•  Meteosat ETv3 reproduces τ in rainfed conditions but is "
        "blind to drip-irrigation pulses (r = 0.07, −70 % bias).\n"
        "•  Implication: satellite ET alone cannot validate or quantify "
        "drip-irrigated water use; it needs a tower-based amplitude "
        "correction."
    )
    add_textbox(slide, M, th_top + Inches(1.1),
                  th_w, bottom_strip_h - Inches(1.1),
                  takehome, font_size=BODY_PT, color=TEXT_DARK,
                  fill=GREY_BG, line=RGBColor(0xCC, 0xCC, 0xCC))

    # =============================================================
    out = Path(__file__).resolve().parent / "poster_template_A0.pptx"
    prs.save(out)
    print(f"[save] {out}")
    print(f"       page = {PAGE_W.inches:.2f}\" × {PAGE_H.inches:.2f}\""
            f"  (A0 portrait, ~84 × 119 cm)")


if __name__ == "__main__":
    build()
